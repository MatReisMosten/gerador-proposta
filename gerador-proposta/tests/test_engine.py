"""Testes das funções puras de generator/engine.py.

Cobre a lógica crítica sem cobertura prévia: classificação de papel de
texto (usada para decidir tamanho/estilo de fonte), validação de tokens
nomeados {TOKEN} e o scanner que lê tokens de um .pptx real.
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from generator.engine import (
    _guess_token_role,
    _normalize_token_key,
    classify_role,
    is_fillable_named_token,
    scan_named_tokens,
)


class TestClassifyRole:
    def test_display_for_large_font(self):
        assert classify_role(72, "Algum texto", slide_no=5, slot=0) == "display"

    def test_cover_title_on_cover_slides(self):
        assert classify_role(44, "Transformação", slide_no=1, slot=0) == "cover_title"

    def test_cover_title_size_ignored_outside_cover_slides(self):
        # Mesmo tamanho, fora dos slides de capa (1,2,19,20) -> cai em "title"
        assert classify_role(44, "Transformação", slide_no=5, slot=0) == "title"

    def test_title_for_medium_font(self):
        assert classify_role(30, "Nosso Diagnóstico", slide_no=5, slot=0) == "title"

    def test_subtitle_for_16pt(self):
        assert classify_role(16, "Complemento do título", slide_no=5, slot=0) == "subtitle"

    def test_label_for_small_short_text(self):
        assert classify_role(12, "Prazo", slide_no=5, slot=0) == "label"

    def test_body_as_fallback(self):
        assert classify_role(None, "Texto corrido qualquer", slide_no=5, slot=0) == "body"

    def test_keep_short_symbols_untouched(self):
        assert classify_role(20, "R$", slide_no=3, slot=1) == "keep"
        assert classify_role(20, "1", slide_no=3, slot=1) == "keep"


class TestNormalizeTokenKey:
    def test_bare_name_gets_wrapped(self):
        assert _normalize_token_key("NOME") == "{NOME}"

    def test_already_wrapped_is_unchanged(self):
        assert _normalize_token_key("{NOME}") == "{NOME}"

    def test_empty_input_returns_empty(self):
        assert _normalize_token_key("") == ""
        assert _normalize_token_key(None) == ""


class TestIsFillableNamedToken:
    @pytest.mark.parametrize(
        "token",
        ["{TITULO_CAPA}", "{BREVE_DESCRICAO}", "TITULO_DOR", "{ITEM_1_DESAFIO}"],
    )
    def test_valid_named_tokens_are_fillable(self, token):
        assert is_fillable_named_token(token) is True

    @pytest.mark.parametrize(
        "token",
        [
            "{}",
            "{ }",
            "{Logo_Cliente}",
            "{LOGO_CLIENTE}",
            "{S01_T00}",
            "{s12_t03}",
            "",
        ],
    )
    def test_non_fillable_tokens_are_rejected(self, token):
        assert is_fillable_named_token(token) is False

    def test_raw_text_with_spaces_is_not_a_token(self):
        # Texto cru do template (rótulos fixos) nunca deve virar token preenchível.
        assert is_fillable_named_token("O DESAFIO") is False


class TestGuessTokenRole:
    @pytest.mark.parametrize(
        "token,expected_role",
        [
            ("{COD_PROJ}", "meta"),
            ("{VALOR_ENTREGA}", "meta"),
            ("{STEP_1}", "step"),
            ("{SEMANAS}", "step"),
            ("{BULLET_1}", "bullet"),
            ("{ITENS_RISCO}", "bullet"),
            ("{SUB_TITULO}", "subtitle"),
            ("{TITULO_CAPA}", "title"),
            ("{RESULT_TITULO}", "title"),
            # NB: mesmo com prefixo CARD_, o sufixo _TITULO cai na regra
            # geral "termina em _TITULO" (checada antes) e vira "title",
            # não "label" — a branch específica de CARD_..._TITULO é
            # inalcançável na ordem atual de _guess_token_role. Ver nota
            # de bug no plano de refatoração.
            ("{CARD_1_TITULO}", "title"),
            ("{IMPACTO_1}", "card_desc"),
            ("{PREM_1}", "card_desc"),
            ("{RESULT_DESC}", "narrative"),
            ("{DESC_ABERTURA}", "narrative"),
        ],
    )
    def test_known_token_families(self, token, expected_role):
        assert _guess_token_role(token, section="") == expected_role

    def test_cover_section_fallback(self):
        assert _guess_token_role("{ALGO_SEM_PADRAO}", section="CAPA") == "cover"

    def test_unknown_token_defaults_to_narrative(self):
        assert _guess_token_role("{ALGO_SEM_PADRAO}", section="OUTRA SEÇÃO") == "narrative"


def _build_sample_pptx(path: Path) -> Path:
    """Cria um .pptx mínimo com 2 slides para testar o scanner de tokens."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    slide1 = prs.slides.add_slide(blank_layout)
    box1 = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box1.text_frame.text = "{TITULO_CAPA}"
    box1.text_frame.paragraphs[0].runs[0].font.size = Pt(40)

    box1b = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    box1b.text_frame.text = "{}"  # placeholder vazio — deve ser ignorado

    slide2 = prs.slides.add_slide(blank_layout)
    box2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box2.text_frame.text = "{TITULO_CAPA}"  # repetido em outro slide

    box2b = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    box2b.text_frame.text = "{S01_T00}"  # lixo legado — deve ser ignorado

    prs.save(str(path))
    return path


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
    return _build_sample_pptx(tmp_path / "sample.pptx")


class TestScanNamedTokens:
    def test_finds_named_token_across_slides(self, sample_pptx):
        found = scan_named_tokens(sample_pptx)
        assert "{TITULO_CAPA}" in found
        assert found["{TITULO_CAPA}"]["count"] == 2
        assert found["{TITULO_CAPA}"]["slides"] == [1, 2]

    def test_ignores_empty_and_legacy_slot_tokens(self, sample_pptx):
        found = scan_named_tokens(sample_pptx)
        assert "{}" not in found
        assert "{S01_T00}" not in found
