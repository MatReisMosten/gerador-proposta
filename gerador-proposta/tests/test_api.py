"""Testes de contrato da API FastAPI (Fase 2) — sem depender de OpenAI real.

Cobre os 3 modos de geração (package / llm_package / llm_full) via
monkeypatch de `fill_slots` para os modos que chamam LLM, e valida os dois
bugs encontrados durante a construção desta API:

1. Escopo Fechado (mode="llm_package") deve isolar a própria seção do
   template, não excluí-la (bug presente no app Streamlit — ver
   api/service.py).
2. Quando nenhum logo é resolvido, o token "{LOGO_CLIENTE}" não pode
   sobrar literal no slide exportado (vazamento presente no app Streamlit
   também).
"""

from __future__ import annotations

import io
import json

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation

from api import config, service
from api.main import app

client = TestClient(app)


def _fake_fill_slots(*, catalog, **_kwargs) -> dict[str, str]:
    return {key: f"Texto de teste para {key}" for key in catalog}


@pytest.fixture(autouse=True)
def _fake_openai_key(monkeypatch):
    # Só precisa existir para passar do guard `if not config.OPENAI_API_KEY`.
    # Os testes que chamam o LLM sempre usam _fake_fill_slots via monkeypatch.
    monkeypatch.setattr(config, "OPENAI_API_KEY", "sk-test-dummy")


class TestHealthAndTypes:
    def test_health(self):
        resp = client.get("/api/health")
        assert resp.status_code == 200
        assert resp.json() == {"ok": True}

    def test_proposal_types_shape(self):
        resp = client.get("/api/proposal-types")
        assert resp.status_code == 200
        by_id = {t["id"]: t for t in resp.json()}

        assert set(by_id) == {
            "professional_service",
            "suporte",
            "passlog",
            "discovery",
            "clarion",
            "escopo_fechado",
            "livre",
        }
        assert by_id["clarion"]["requires_form"] is False
        assert by_id["discovery"]["show_brief"] is True
        assert by_id["professional_service"]["show_brief"] is False
        assert by_id["escopo_fechado"]["hide_client"] is True
        assert by_id["escopo_fechado"]["hide_logo"] is True
        assert by_id["livre"]["show_attachments"] is True

    def test_templates_summary(self):
        resp = client.get("/api/templates/summary")
        assert resp.status_code == 200
        data = resp.json()
        assert data["name"] == "slide-mestre-template.pptx"
        assert "Escopo Fechado (DP World)" in data["sections"]


class TestValidation:
    def test_unknown_type_is_404(self):
        resp = client.post("/api/proposals/generate", data={"type_id": "nao_existe"})
        assert resp.status_code == 404

    def test_invalid_project_code_is_422(self):
        resp = client.post(
            "/api/proposals/generate",
            data={"type_id": "suporte", "project_code": "ABC"},
        )
        assert resp.status_code == 422
        assert "AAA999-99" in resp.json()["detail"]

    def test_missing_required_fields_is_422(self):
        resp = client.post(
            "/api/proposals/generate",
            data={
                "type_id": "suporte",
                "project_code": "SUP001-26",
                "fields_json": "{}",
            },
        )
        assert resp.status_code == 422
        assert "obrigat" in resp.json()["detail"].lower()

    def test_invalid_fields_json_is_422(self):
        resp = client.post(
            "/api/proposals/generate",
            data={
                "type_id": "suporte",
                "project_code": "SUP001-26",
                "fields_json": "[not json}",
            },
        )
        assert resp.status_code == 422

    def test_download_unknown_token_is_404(self):
        resp = client.get("/api/proposals/download/does-not-exist")
        assert resp.status_code == 404


class TestPackageGeneration:
    def test_professional_service_generates_and_downloads(self):
        resp = client.post(
            "/api/proposals/generate",
            data={
                "type_id": "professional_service",
                "client_name": "Cliente Teste",
                "project_code": "TST001-26",
                "fields_json": json.dumps(
                    {
                        "total": "R$ 28.800,00",
                        "tempo_execucao": "3",
                        "tipo_profissional": "Desenvolvedor QA",
                        "senioridade": "Sênior",
                    }
                ),
            },
        )
        assert resp.status_code == 200
        body = resp.json()
        assert body["meta"]["code"] == "TST001-26"
        assert body["meta"]["filename"].endswith(".pptx")

        dl = client.get(f"/api/proposals/download/{body['token']}")
        assert dl.status_code == 200
        assert dl.headers["content-type"].startswith("application/vnd")

        prs = Presentation(io.BytesIO(dl.content))
        texts = [
            shape.text_frame.text
            for slide in prs.slides
            for shape in slide.shapes
            if shape.has_text_frame
        ]
        assert not any("{LOGO_CLIENTE}" in t for t in texts), (
            "token de logo vazou como texto literal quando nenhum logo foi enviado"
        )

    def test_clarion_defaults_invalid_code(self):
        resp = client.post("/api/proposals/generate", data={"type_id": "clarion"})
        assert resp.status_code == 200
        assert resp.json()["meta"]["code"] == "CLA000-00"


class TestLlmPackageGeneration:
    """Cobre o bug fix: Escopo Fechado isola a própria seção (não exclui)."""

    def test_escopo_fechado_isolates_its_own_section(self, monkeypatch):
        monkeypatch.setattr(service, "fill_slots", _fake_fill_slots)

        resp = client.post(
            "/api/proposals/generate",
            data={
                "type_id": "escopo_fechado",
                "project_code": "DPW001-26",
                "brief": "Contexto de teste para o Escopo Fechado.",
            },
        )
        assert resp.status_code == 200
        body = resp.json()
        assert body["meta"]["code"] == "DPW001-26"
        assert body["meta"]["client"] == "", "hide_client deve zerar o nome do cliente"

        summary = client.get("/api/templates/summary").json()
        expected_slides = summary["sections"]["Escopo Fechado (DP World)"]

        dl = client.get(f"/api/proposals/download/{body['token']}")
        prs = Presentation(io.BytesIO(dl.content))
        assert len(prs.slides) == expected_slides

    def test_escopo_fechado_without_brief_is_422(self):
        resp = client.post(
            "/api/proposals/generate",
            data={"type_id": "escopo_fechado", "project_code": "DPW001-26"},
        )
        assert resp.status_code == 422


class TestLivreGeneration:
    def test_livre_uses_full_catalog(self, monkeypatch):
        monkeypatch.setattr(service, "fill_slots", _fake_fill_slots)

        resp = client.post(
            "/api/proposals/generate",
            data={
                "type_id": "livre",
                "client_name": "Cliente Livre",
                "project_code": "LIV001-26",
                "brief": "Brief de teste para o modo Livre.",
            },
        )
        assert resp.status_code == 200
        assert resp.json()["meta"]["type"] == "livre"

    def test_livre_without_brief_is_422(self):
        resp = client.post(
            "/api/proposals/generate",
            data={"type_id": "livre", "project_code": "LIV001-26"},
        )
        assert resp.status_code == 422
