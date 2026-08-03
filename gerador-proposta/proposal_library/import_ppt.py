"""Orchestrate PPTX import into a staging directory (no AI)."""

from __future__ import annotations

import logging
import uuid
from datetime import datetime, timezone
from pathlib import Path

from pptx import Presentation

from .extractor import extract_slide
from .image_extractor import extract_images
from .importer import isolate_slide
from .preview_generator import PreviewError, generate_preview
from .utils import ensure_dir, resolve_path, slide_stem, write_json

logger = logging.getLogger(__name__)


def _batch_id() -> str:
    stamp = datetime.now(timezone.utc).strftime("%Y%m%dT%H%M%SZ")
    return f"{stamp}_{uuid.uuid4().hex[:8]}"


def import_ppt(
    pptx_path: Path,
    out_dir: Path | None = None,
    *,
    skip_preview: bool = False,
) -> Path:
    """
    Split a PPTX into per-slide artifacts under out_dir.

    Layout:
      out_dir/
        manifest.json
        slide_001.pptx
        slide_001.png
        slide_001.json
        images/slide_001/image_1.png
        enrichment/   (empty, for the agent)
    """
    pptx_path = resolve_path(pptx_path)
    if not pptx_path.is_file():
        raise FileNotFoundError(f"PPTX não encontrado: {pptx_path}")

    if out_dir is None:
        raise ValueError("out_dir é obrigatório")

    out_dir = resolve_path(out_dir)
    ensure_dir(out_dir)
    ensure_dir(out_dir / "enrichment")
    ensure_dir(out_dir / "images")

    prs = Presentation(str(pptx_path))
    total = len(prs.slides)
    if total == 0:
        raise ValueError(f"Apresentação sem slides: {pptx_path}")

    logger.info("Importing %d slide(s) from %s", total, pptx_path.name)
    slides_meta: list[dict] = []

    for index in range(total):
        stem = slide_stem(index + 1)
        slide_pptx = out_dir / f"{stem}.pptx"
        slide_png = out_dir / f"{stem}.png"
        slide_json = out_dir / f"{stem}.json"
        images_dir = out_dir / "images" / stem

        logger.info("Isolating slide %d/%d", index + 1, total)
        isolate_slide(pptx_path, index, slide_pptx)

        image_names = extract_images(slide_pptx, images_dir)

        # Re-open isolated deck for extraction (single slide at index 0)
        isolated = Presentation(str(slide_pptx))
        extraction = extract_slide(
            isolated,
            0,
            image_names=image_names,
            source_file=pptx_path.name,
        )
        # Keep original deck slide number for humans
        extraction.slide = index + 1

        if skip_preview:
            logger.warning("skip_preview=True — gerando PNG placeholder para %s", stem)
            from PIL import Image

            Image.new("RGB", (1280, 720), color=(240, 240, 240)).save(slide_png)
            extraction.visual_hash = ""
        else:
            try:
                extraction.visual_hash = generate_preview(slide_pptx, slide_png)
            except PreviewError:
                # Clean partial artifacts? Keep pptx/json/images; re-raise.
                raise

        write_json(slide_json, extraction.to_dict())
        slides_meta.append(
            {
                "slide": extraction.slide,
                "stem": stem,
                "pptx": slide_pptx.name,
                "png": slide_png.name if slide_png.exists() else None,
                "json": slide_json.name,
                "images_dir": f"images/{stem}",
                "image_count": len(image_names),
            }
        )
        logger.info(
            "Slide %d ready (%d texts, %d images)",
            extraction.slide,
            len(extraction.texts),
            len(image_names),
        )

    manifest = {
        "batch_id": out_dir.name,
        "source": str(pptx_path),
        "source_name": pptx_path.name,
        "created_at": datetime.now(timezone.utc).isoformat(),
        "slide_count": total,
        "slides": slides_meta,
        "enrichment_dir": "enrichment",
    }
    write_json(out_dir / "manifest.json", manifest)
    logger.info("Staging ready: %s (%d slides)", out_dir, total)
    return out_dir


def default_staging_dir(library_root: Path) -> Path:
    return library_root / "_staging" / _batch_id()
