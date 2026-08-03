"""Apply placeholder text replacements inside a PPTX slide."""

from __future__ import annotations

import logging
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .models import PlaceholderMapping

logger = logging.getLogger(__name__)


def _shape_full_text(shape) -> str:
    if not hasattr(shape, "text_frame"):
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs).strip()


def _iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)


def _set_shape_text(shape, new_text: str) -> None:
    """Replace shape text while keeping the first paragraph/run formatting."""
    tf = shape.text_frame
    for para in list(tf.paragraphs[1:]):
        para._p.getparent().remove(para._p)  # noqa: SLF001

    # Support multi-line placeholders
    lines = new_text.split("\n")
    first = tf.paragraphs[0]
    for run in first.runs:
        run.text = ""
    if first.runs:
        first.runs[0].text = lines[0] if lines else new_text
    else:
        run = first.add_run()
        run.text = lines[0] if lines else new_text

    for line in lines[1:]:
        para = tf.add_paragraph()
        run = para.add_run()
        run.text = line


def replace_in_shape(shape, replacements: list[tuple[str, str]]) -> bool:
    """
    Apply replacements to a shape.

    Prefers exact full-text match; falls back to substring replace in runs.
    """
    if not hasattr(shape, "text_frame"):
        return False

    full = _shape_full_text(shape)
    if not full:
        return False

    # Exact / contained full-shape replacements (longest first)
    ordered = sorted(replacements, key=lambda x: len(x[0]), reverse=True)
    for original, placeholder in ordered:
        if not original:
            continue
        if full == original or full.replace("\xa0", " ") == original:
            _set_shape_text(shape, placeholder)
            return True

    changed = False
    for original, placeholder in ordered:
        if not original or original not in full:
            continue
        # Replace inside runs when possible
        remaining = original
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if remaining and remaining in run.text:
                    run.text = run.text.replace(remaining, placeholder)
                    changed = True
                    remaining = ""
                elif remaining and remaining[: len(run.text)] == run.text and run.text:
                    # Multi-run span: clear middle runs after first partial — conservative
                    pass
        if changed:
            # Re-read and do a full-shape replace if still contains original
            full2 = _shape_full_text(shape)
            if original in full2:
                _set_shape_text(shape, full2.replace(original, placeholder))
            return True

    return changed


def apply_placeholders(pptx_path: Path, mappings: list[PlaceholderMapping]) -> int:
    """
    Apply placeholder mappings to all slides in the PPTX.
    Returns number of shapes modified.
    """
    if not mappings:
        return 0

    replacements = [(m.original, m.placeholder) for m in mappings if m.original]
    prs = Presentation(str(pptx_path))
    modified = 0

    for slide in prs.slides:
        for shape in _iter_shapes(slide.shapes):
            if replace_in_shape(shape, replacements):
                modified += 1
                logger.debug("Replaced text in shape %s", getattr(shape, "name", "?"))

    prs.save(str(pptx_path))
    return modified


def collect_texts_after_placeholders(pptx_path: Path) -> list[str]:
    prs = Presentation(str(pptx_path))
    texts: list[str] = []
    for slide in prs.slides:
        for shape in _iter_shapes(slide.shapes):
            text = _shape_full_text(shape)
            if text:
                texts.append(text)
    return texts
