"""Assemble a multi-slide PPTX from library slide templates + placeholder values."""

from __future__ import annotations

import logging
import shutil
import tempfile
from copy import deepcopy
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from .placeholder_applier import apply_placeholders
from .models import PlaceholderMapping
from .utils import ensure_dir, resolve_path

logger = logging.getLogger(__name__)

_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"


def _fill_slide(src_pptx: Path, values: dict[str, str], dest_pptx: Path) -> None:
    ensure_dir(dest_pptx.parent)
    shutil.copy2(src_pptx, dest_pptx)
    mappings = [
        PlaceholderMapping(original=ph, placeholder=text)
        for ph, text in values.items()
        if ph and text is not None
    ]
    # apply_placeholders replaces original→placeholder; we need placeholder→value
    # So invert: treat current {{X}} as "original" and value as "placeholder" (new text)
    inverted = [
        PlaceholderMapping(original=ph, placeholder=val)
        for ph, val in values.items()
    ]
    apply_placeholders(dest_pptx, inverted)


def _iter_pictures(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_pictures(shape.shapes)
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            yield shape


def _replace_picture_blob(shape, image_bytes: bytes) -> bool:
    try:
        # python-pptx oxml elements don't accept namespaces= in xpath
        blips = shape._element.findall(  # noqa: SLF001
            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}blip"
        )
        if not blips:
            return False
        embed = blips[0].get(qn("r:embed"))
        if not embed:
            return False
        image_part = shape.part.related_part(embed)
        image_part.blob = image_bytes
        return True
    except Exception as exc:
        logger.warning("Could not replace picture blob: %s", exc)
        return False


def replace_client_logo(
    pptx_path: Path,
    new_logo: Path,
    *,
    reference_blob: bytes | None = None,
) -> bool:
    """
    Replace client logo image in a slide.

    If reference_blob is provided (original IMG_LOGO_CLIENT bytes), match by content.
    Otherwise replace the smallest non-icon picture (typical logo slot).
    """
    prs = Presentation(str(pptx_path))
    if not prs.slides:
        return False

    new_bytes = new_logo.read_bytes()
    pictures = list(_iter_pictures(prs.slides[0].shapes))
    if not pictures:
        return False

    target = None
    if reference_blob:
        for shape in pictures:
            try:
                if shape.image.blob == reference_blob:
                    target = shape
                    break
            except Exception:
                continue

    if target is None:
        slide_area = (prs.slide_width or 1) * (prs.slide_height or 1)
        ranked = sorted(
            pictures,
            key=lambda s: (s.width or 0) * (s.height or 0),
        )
        # Prefer small-to-medium pictures (logos), skip tiny icons and huge heroes
        for shape in ranked:
            area = (shape.width or 0) * (shape.height or 0)
            if slide_area * 0.0005 < area < slide_area * 0.12:
                target = shape
                break
        if target is None and ranked:
            target = ranked[0]

    if target is None:
        return False

    ok = _replace_picture_blob(target, new_bytes)
    if ok:
        prs.save(str(pptx_path))
    return ok


def _append_slide(dest: Presentation, source_pptx: Path) -> None:
    """Append the first slide of source_pptx into dest, preserving layout XML."""
    source = Presentation(str(source_pptx))
    if not source.slides:
        raise ValueError(f"No slides in {source_pptx}")

    src_slide = source.slides[0]
    # Use blank layout as carrier; then replace tree with deep copy of source slide
    blank = dest.slide_layouts[6] if len(dest.slide_layouts) > 6 else dest.slide_layouts[0]
    dest_slide = dest.slides.add_slide(blank)

    # Remove default shapes from blank slide
    for shape in list(dest_slide.shapes):
        sp = shape._element  # noqa: SLF001
        sp.getparent().remove(sp)

    # Deep-copy shape tree
    src_tree = src_slide.shapes._spTree  # noqa: SLF001
    dest_tree = dest_slide.shapes._spTree  # noqa: SLF001

    # Copy non-nvGrpSpPr children (actual shapes)
    for child in list(src_tree):
        tag = etree.QName(child.tag).localname
        if tag in {"nvGrpSpPr", "grpSpPr"}:
            continue
        dest_tree.append(deepcopy(child))

    # Copy image parts / relationships referenced by copied blips
    for blip in dest_tree.xpath(
        ".//a:blip",
        namespaces={"a": "http://schemas.openxmlformats.org/drawingml/2006/main"},
    ):
        embed = blip.get(qn("r:embed"))
        if not embed:
            continue
        try:
            # Relationship may still point to source — re-wire
            # Find matching rId in source slide part
            src_rel = src_slide.part.rels[embed]
            image_blob = src_rel.target_part.blob
            content_type = src_rel.target_part.content_type
            # Add image to dest slide part
            new_rid, _ = dest_slide.part.get_or_add_image(
                # get_or_add_image expects a file-like or we use related parts API
            )
        except Exception:
            pass

    # Simpler reliable path: use package-level approach below instead


def merge_single_slide_decks(slide_paths: list[Path], output: Path) -> None:
    """
    Merge single-slide PPTX files into one deck.

    Strategy: start from first file; for each next file, add a blank slide and
    deep-copy shapes, rewiring image relationships.
    """
    from io import BytesIO

    from PIL import Image

    if not slide_paths:
        raise ValueError("No slides to merge")

    ensure_dir(output.parent)
    shutil.copy2(slide_paths[0], output)

    if len(slide_paths) == 1:
        return

    def _image_stream(blob: bytes) -> BytesIO:
        try:
            stream = BytesIO(blob)
            img = Image.open(stream)
            out = BytesIO()
            img.convert("RGBA").save(out, format="PNG")
            out.seek(0)
            return out
        except Exception:
            out = BytesIO()
            Image.new("RGBA", (1, 1), (0, 0, 0, 0)).save(out, format="PNG")
            out.seek(0)
            return out

    dest = Presentation(str(output))

    for path in slide_paths[1:]:
        src = Presentation(str(path))
        src_slide = src.slides[0]

        layout = dest.slide_layouts[6] if len(dest.slide_layouts) > 6 else dest.slide_layouts[0]
        blank_slide = dest.slides.add_slide(layout)

        for el in list(blank_slide.shapes._spTree):  # noqa: SLF001
            if any(
                el.tag.endswith(suffix)
                for suffix in ("}sp", "}pic", "}grpSp", "}cxnSp", "}graphicFrame")
            ):
                blank_slide.shapes._spTree.remove(el)  # noqa: SLF001

        rid_map: dict[str, str] = {}
        for rel in src_slide.part.rels.values():
            if "image" not in (rel.reltype or ""):
                continue
            try:
                stream = _image_stream(rel.target_part.blob)
                _part, new_rid = blank_slide.part.get_or_add_image_part(stream)
                rid_map[rel.rId] = new_rid
            except Exception as exc:
                logger.warning("Skipping image %s from %s: %s", rel.rId, path.name, exc)

        for child in src_slide.shapes._spTree:  # noqa: SLF001
            tag = etree.QName(child.tag).localname
            if tag in {"nvGrpSpPr", "grpSpPr"}:
                continue
            new_child = deepcopy(child)
            for blip in new_child.findall(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}blip"
            ):
                old = blip.get(qn("r:embed"))
                if old and old in rid_map:
                    blip.set(qn("r:embed"), rid_map[old])
            blank_slide.shapes._spTree.append(new_child)  # noqa: SLF001

    dest.save(str(output))
    logger.info("Merged %d slides -> %s", len(slide_paths), output)


def generate_proposal(
    *,
    library_root: Path,
    slide_plan: list[dict],
    output_pptx: Path,
    logo_path: Path | None = None,
    work_dir: Path | None = None,
) -> Path:
    """
    slide_plan: [{"id": "cover_003", "values": {"{{COVER_TITLE}}": "..."}}, ...]
    """
    library_root = resolve_path(library_root)
    output_pptx = resolve_path(output_pptx)

    tmp_root = Path(work_dir) if work_dir else Path(tempfile.mkdtemp(prefix="plb_gen_"))
    ensure_dir(tmp_root)
    filled: list[Path] = []

    for i, item in enumerate(slide_plan, start=1):
        slide_id = item["id"]
        values = item.get("values", {})
        src = library_root / "slides" / slide_id / "slide.pptx"
        if not src.is_file():
            raise FileNotFoundError(f"Library slide not found: {slide_id}")
        dest = tmp_root / f"{i:02d}_{slide_id}.pptx"
        _fill_slide(src, values, dest)
        if logo_path and item.get("replace_logo") and Path(logo_path).is_file():
            ref_path = library_root / "slides" / slide_id / "images" / "IMG_LOGO_CLIENT.png"
            ref_blob = ref_path.read_bytes() if ref_path.is_file() else None
            ok = replace_client_logo(dest, Path(logo_path), reference_blob=ref_blob)
            logger.info("Logo replace on %s: %s", slide_id, ok)
        filled.append(dest)

    merge_single_slide_decks(filled, output_pptx)
    return output_pptx
