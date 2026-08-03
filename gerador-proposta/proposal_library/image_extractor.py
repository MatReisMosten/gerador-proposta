"""Extract embedded images from a single-slide PPTX."""

from __future__ import annotations

import logging
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .utils import ensure_dir

logger = logging.getLogger(__name__)

_CONTENT_TYPE_EXT = {
    "image/png": ".png",
    "image/jpeg": ".jpg",
    "image/jpg": ".jpg",
    "image/gif": ".gif",
    "image/bmp": ".bmp",
    "image/tiff": ".tiff",
    "image/x-emf": ".emf",
    "image/x-wmf": ".wmf",
    "image/svg+xml": ".svg",
}


def _iter_picture_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_picture_shapes(shape.shapes)
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            yield shape


def _ext_for_image(image) -> str:
    content_type = getattr(image, "content_type", "") or ""
    if content_type in _CONTENT_TYPE_EXT:
        return _CONTENT_TYPE_EXT[content_type]
    ext = getattr(image, "ext", None)
    if ext:
        return f".{ext.lstrip('.')}"
    return ".bin"


def extract_images(pptx_path: Path, images_dir: Path) -> list[str]:
    """
    Save picture shapes from the first (and typically only) slide.

    Returns relative filenames like image_1.png.
    """
    ensure_dir(images_dir)
    prs = Presentation(str(pptx_path))
    if not prs.slides:
        return []

    saved: list[str] = []
    seen_blobs: set[bytes] = set()
    counter = 0

    for shape in _iter_picture_shapes(prs.slides[0].shapes):
        try:
            image = shape.image
            blob = image.blob
        except Exception as exc:
            logger.warning("Skipping unreadable picture in %s: %s", pptx_path.name, exc)
            continue

        if blob in seen_blobs:
            continue
        seen_blobs.add(blob)

        counter += 1
        filename = f"image_{counter}{_ext_for_image(image)}"
        out_path = images_dir / filename
        out_path.write_bytes(blob)
        saved.append(filename)
        logger.debug("Extracted %s (%d bytes)", filename, len(blob))

    return saved


def rename_images(
    images_dir: Path,
    labels: dict[str, str],
) -> list[str]:
    """
    Rename image files according to labels mapping:
    {"image_1.png": "IMG_HERO"} -> IMG_HERO.png
    Returns ordered list of final filenames (without path).
    """
    if not images_dir.exists():
        return []

    final_names: list[str] = []
    used: set[str] = set()

    for src_name in sorted(p.name for p in images_dir.iterdir() if p.is_file()):
        src = images_dir / src_name
        label = labels.get(src_name)
        if label:
            stem = label if not label.lower().startswith("img_") else label
            if not stem.upper().startswith("IMG_"):
                stem = f"IMG_{stem}"
            ext = src.suffix or ".png"
            dest_name = f"{stem}{ext}"
        else:
            dest_name = src_name

        # Avoid collisions
        base_stem = Path(dest_name).stem
        ext = Path(dest_name).suffix
        candidate = dest_name
        n = 2
        while candidate.lower() in used or (
            (images_dir / candidate).exists() and candidate != src_name
        ):
            candidate = f"{base_stem}_{n}{ext}"
            n += 1

        if candidate != src_name:
            src.rename(images_dir / candidate)
        used.add(candidate.lower())
        final_names.append(candidate)

    return final_names
