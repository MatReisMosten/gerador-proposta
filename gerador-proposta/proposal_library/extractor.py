"""Extract texts, fonts, colors, positions, notes, and layout from a slide."""

from __future__ import annotations

import logging
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

from .models import ShapeInfo, SlideExtraction
from .utils import normalize_text, sha256_text

logger = logging.getLogger(__name__)


def _emu_to_inches(value: int | None) -> float:
    if value is None:
        return 0.0
    return round(Emu(value).inches, 4)


def _font_size_pt(run) -> float | None:
    try:
        if run.font.size is not None:
            return round(run.font.size.pt, 2)
    except Exception:
        return None
    return None


def _font_name(run) -> str | None:
    try:
        return run.font.name
    except Exception:
        return None


def _font_color_hex(run) -> str | None:
    try:
        color = run.font.color
        if color is None or color.type is None:
            return None
        rgb = color.rgb
        if rgb is None:
            return None
        return f"#{str(rgb)}"
    except Exception:
        return None


def _shape_type_name(shape) -> str:
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return "picture"
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            return "group"
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return "table"
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            return "chart"
        if hasattr(shape, "text_frame"):
            return "textbox"
        return str(shape.shape_type)
    except Exception:
        return "unknown"


def _iter_shapes(shapes) -> list[Any]:
    result: list[Any] = []
    for shape in shapes:
        result.append(shape)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            result.extend(_iter_shapes(shape.shapes))
    return result


def _shape_full_text(shape) -> str:
    if not hasattr(shape, "text_frame"):
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs).strip()


def _collect_run_styles(shape) -> tuple[float | None, str | None, str | None]:
    if not hasattr(shape, "text_frame"):
        return None, None, None
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            size = _font_size_pt(run)
            name = _font_name(run)
            color = _font_color_hex(run)
            if size or name or color:
                return size, name, color
    return None, None, None


def _notes_text(slide) -> str:
    try:
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame is not None:
                return notes_frame.text.strip()
    except Exception as exc:
        logger.debug("Could not read notes: %s", exc)
    return ""


def _layout_label(prs: Presentation) -> str:
    w = prs.slide_width
    h = prs.slide_height
    if not w or not h:
        return "unknown"
    ratio = w / h
    if abs(ratio - 16 / 9) < 0.05:
        return "16:9"
    if abs(ratio - 4 / 3) < 0.05:
        return "4:3"
    return f"{_emu_to_inches(w)}x{_emu_to_inches(h)}"


def _structure_fingerprint(shapes: list[ShapeInfo], layout: str) -> str:
    parts: list[str] = [layout]
    for s in shapes:
        parts.append(
            f"{s.type}:{round(s.x, 2)}:{round(s.y, 2)}:{round(s.w, 2)}:{round(s.h, 2)}"
        )
    return sha256_text("|".join(parts))


def extract_slide(
    prs: Presentation,
    slide_index: int,
    *,
    image_names: list[str] | None = None,
    source_file: str = "",
) -> SlideExtraction:
    """Extract deterministic metadata from a 0-based slide index."""
    slide = prs.slides[slide_index]
    shape_infos: list[ShapeInfo] = []
    texts: list[str] = []
    font_sizes: list[float] = []
    fonts: list[str] = []
    colors: list[str] = []

    for shape in _iter_shapes(slide.shapes):
        text = _shape_full_text(shape)
        font_size, font_name, color = _collect_run_styles(shape)
        info = ShapeInfo(
            name=getattr(shape, "name", "") or "",
            type=_shape_type_name(shape),
            x=_emu_to_inches(getattr(shape, "left", None)),
            y=_emu_to_inches(getattr(shape, "top", None)),
            w=_emu_to_inches(getattr(shape, "width", None)),
            h=_emu_to_inches(getattr(shape, "height", None)),
            text=text,
            font_size=font_size,
            font_name=font_name,
            color=color,
        )
        shape_infos.append(info)

        if text:
            texts.append(text)
            if font_size is not None:
                font_sizes.append(font_size)
            if font_name and font_name not in fonts:
                fonts.append(font_name)
            if color and color not in colors:
                colors.append(color)

    layout = _layout_label(prs)
    title = texts[0] if texts else ""
    notes = _notes_text(slide)

    return SlideExtraction(
        slide=slide_index + 1,
        title=title,
        texts=texts,
        font_sizes=font_sizes,
        fonts=fonts,
        colors=colors,
        shapes=[s.to_dict() for s in shape_infos],
        images=list(image_names or []),
        notes=notes,
        layout=layout,
        element_count=len(shape_infos),
        structure_hash=_structure_fingerprint(shape_infos, layout),
        source_file=source_file,
    )


def texts_for_similarity(extraction: SlideExtraction) -> set[str]:
    return {normalize_text(t) for t in extraction.texts if normalize_text(t)}
