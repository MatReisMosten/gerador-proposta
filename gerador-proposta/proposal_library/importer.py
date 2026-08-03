"""Open PPTX and isolate each slide into its own single-slide PPTX."""

from __future__ import annotations

import logging
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

from .utils import ensure_dir

logger = logging.getLogger(__name__)


def count_slides(pptx_path: Path) -> int:
    prs = Presentation(str(pptx_path))
    return len(prs.slides)


def _sld_id_list(prs: Presentation):
    return prs.slides._sldIdLst  # noqa: SLF001 — required API


def delete_slide(prs: Presentation, index: int) -> None:
    """Remove slide at index from the presentation (in-memory)."""
    sld_id_lst = _sld_id_list(prs)
    sld_ids = list(sld_id_lst)
    if index < 0 or index >= len(sld_ids):
        raise IndexError(f"Slide index out of range: {index}")

    sld_id = sld_ids[index]
    r_id = sld_id.get(qn("r:id"))
    sld_id_lst.remove(sld_id)
    if r_id:
        try:
            prs.part.drop_rel(r_id)
        except KeyError:
            logger.debug("Relationship %s already dropped", r_id)


def isolate_slide(source_pptx: Path, slide_index: int, dest_pptx: Path) -> None:
    """
    Create a PPTX containing only the slide at slide_index (0-based).

    Copies the source file then deletes all other slides.
    """
    ensure_dir(dest_pptx.parent)
    shutil.copy2(source_pptx, dest_pptx)
    prs = Presentation(str(dest_pptx))
    total = len(prs.slides)
    if slide_index < 0 or slide_index >= total:
        dest_pptx.unlink(missing_ok=True)
        raise IndexError(f"Slide {slide_index} not in range 0..{total - 1}")

    # Delete from the end to keep indices stable until we hit the kept slide.
    for i in range(total - 1, -1, -1):
        if i != slide_index:
            delete_slide(prs, i)

    prs.save(str(dest_pptx))
    logger.debug("Isolated slide %d -> %s", slide_index + 1, dest_pptx.name)
