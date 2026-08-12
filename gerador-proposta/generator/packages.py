"""Proposal type packages: registry, section isolation, placeholder-only fill."""

from __future__ import annotations

import json
import re
import shutil
from datetime import date
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from . import paths as P
from .engine import apply_named_placeholders, keep_only_slides, replace_logo_cliente

_PLACEHOLDER_RE = re.compile(r"\{([^{}]+)\}")


def load_packages_registry() -> dict[str, Any]:
    path = P.packages_registry_path()
    if not path.is_file():
        return {"version": 1, "types": [{"id": "livre", "label": "Livre", "mode": "llm_full"}]}
    return json.loads(path.read_text(encoding="utf-8"))


def list_proposal_types() -> list[dict[str, Any]]:
    return list(load_packages_registry().get("types") or [])


def get_proposal_type(type_id: str) -> dict[str, Any]:
    for item in list_proposal_types():
        if item.get("id") == type_id:
            return item
    raise KeyError(f"Tipo de proposta desconhecido: {type_id}")


def read_pptx_sections(pptx_path: Path) -> dict[str, list[int]]:
    """
    Read PowerPoint section names → list of 1-based slide indices.
    """
    prs = Presentation(str(pptx_path))
    sld_ids = list(prs.slides._sldIdLst)  # noqa: SLF001
    id_to_index: dict[str, int] = {}
    for i, sld in enumerate(sld_ids, start=1):
        sid = sld.get("id")
        if sid:
            id_to_index[sid] = i

    sections: dict[str, list[int]] = {}
    root = prs.part._element  # noqa: SLF001
    for el in root.iter():
        if etree.QName(el).localname != "section":
            continue
        name = (el.get("name") or "").strip()
        if not name:
            continue
        slides: list[int] = []
        for child in el.iter():
            if etree.QName(child).localname != "sldId":
                continue
            sid = child.get("id")
            if sid and sid in id_to_index:
                slides.append(id_to_index[sid])
        sections[name] = slides
    return sections


def resolve_package_slides(pkg: dict[str, Any], pptx_path: Path) -> list[int]:
    """Resolve which 1-based slides belong to this package."""
    explicit = pkg.get("slides_1based")
    if explicit:
        return [int(x) for x in explicit]

    section_name = (pkg.get("section") or "").strip()
    if section_name:
        sections = read_pptx_sections(pptx_path)
        if section_name not in sections:
            available = ", ".join(sorted(sections)) or "(nenhuma)"
            raise KeyError(
                f"Seção '{section_name}' não encontrada no template. "
                f"Disponíveis: {available}"
            )
        slides = sections[section_name]
        if not slides:
            raise ValueError(f"Seção '{section_name}' não contém slides.")
        return slides

    # whole deck
    prs = Presentation(str(pptx_path))
    return list(range(1, len(prs.slides) + 1))


def resolve_package_template(pkg: dict[str, Any]) -> Path:
    """
    Resolve PPTX source for package-mode types.
    Professional Service / source=master → sempre slide-mestre-template.pptx.
    """
    source = (pkg.get("source") or "standalone").lower()
    rel = (pkg.get("template") or "").strip()

    # Master / mestre: never the legacy Modelo-Proposta-Tecnica
    if source == "master" or rel.endswith(P.MASTER_TEMPLATE_NAME) or "mestre" in rel.lower():
        path = P.master_template_path()
        if path.name != P.MASTER_TEMPLATE_NAME:
            raise FileNotFoundError(
                f"Template master inválido: {path}. "
                f"Esperado: {P.MASTER_TEMPLATE_NAME}"
            )
        return path

    candidates = [
        P.APP_ROOT / "data" / rel,
        P.UNISANTA / rel,
        P.packages_dir() / Path(rel).name,
    ]
    for path in candidates:
        if path and path.is_file():
            # Guard: never silently use the old parameterized deck for packages
            if "Modelo-Proposta-Tecnica" in path.name:
                continue
            return path.resolve()

    raise FileNotFoundError(
        f"Template do pacote '{pkg.get('id')}' não encontrado. "
        f"Esperado: data/{P.MASTER_TEMPLATE_NAME} (source=master) "
        f"ou {rel}"
    )


def build_field_values(
    pkg: dict[str, Any],
    *,
    field_values: dict[str, str],
    client_name: str = "",
    project_code: str = "",
) -> dict[str, str]:
    """Map package placeholder keys → concrete strings (only named {TOKENS})."""
    fields = {k: (v or "").strip() for k, v in field_values.items()}
    code = (
        fields.get("codigo_projeto")
        or fields.get("project_code")
        or project_code.strip()
        or "A definir"
    )
    client = (
        fields.get("nome_cliente")
        or fields.get("client_name")
        or client_name.strip()
        or "Cliente"
    )
    context = {
        "client_name": client,
        "nome_cliente": client,
        "project_code": code,
        "codigo_projeto": code,
        "today": date.today().strftime("%d/%m/%Y"),
        **fields,
    }
    # Keep shared aliases even when package fields override
    context["client_name"] = client
    context["nome_cliente"] = client
    context["project_code"] = code
    context["codigo_projeto"] = code

    # {DESC_SEMANAS}: "8 semanas" a partir do número informado
    tempo = (fields.get("tempo_execucao") or "").strip()
    if tempo:
        if "semana" in tempo.lower():
            context["desc_semanas"] = tempo
        else:
            context["desc_semanas"] = f"{tempo} semanas"
    else:
        context.setdefault("desc_semanas", "")

    mapping = pkg.get("placeholders") or {}
    result: dict[str, str] = {}
    for ph, source_key in mapping.items():
        token = ph.strip()
        if not token or token == "{}":
            continue
        # never overwrite logo markers — handled by replace_logo_cliente
        if token in P.KEEP_AS_IS or token.strip("{}").upper() == "LOGO_CLIENTE":
            continue
        result[token] = str(context.get(source_key, "") or "")
    return result


def _iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)


def fill_investment_table(
    pptx_path: Path,
    *,
    perfil: str = "",
    quantidade: str = "",
    valor_hora: str = "",
    total: str = "",
) -> bool:
    """
    Fill the first blank data row of the 'Modelo de Investimento' table.
    Does not alter header labels (texto bruto).
    """
    if not any([perfil, quantidade, valor_hora, total]):
        return False

    prs = Presentation(str(pptx_path))
    changed = False
    for slide in prs.slides:
        titles = []
        for shape in _iter_shapes(slide.shapes):
            if shape.has_text_frame and shape.text_frame.text.strip():
                titles.append(shape.text_frame.text.strip().lower())
        if not any("modelo de investimento" in t for t in titles):
            continue
        for shape in _iter_shapes(slide.shapes):
            if shape.shape_type != MSO_SHAPE_TYPE.TABLE:
                continue
            table = shape.table
            # row0 = headers; fill first empty data row; then Total row
            filled_data = False
            for ri in range(1, len(table.rows)):
                cells = [c.text.strip() for c in table.rows[ri].cells]
                if cells and cells[0].lower() == "total":
                    if total:
                        # put total on the rightmost empty value cell
                        for ci in range(len(table.rows[ri].cells) - 1, 0, -1):
                            if not table.rows[ri].cells[ci].text.strip():
                                table.rows[ri].cells[ci].text = total
                                changed = True
                                break
                    continue
                if filled_data or any(cells):
                    continue
                vals = [perfil, quantidade, valor_hora]
                for ci, val in enumerate(vals):
                    if ci < len(table.rows[ri].cells) and val:
                        table.rows[ri].cells[ci].text = val
                        changed = True
                filled_data = True
        if changed:
            break
    if changed:
        prs.save(str(pptx_path))
    return changed


def build_package_deck(
    pkg: dict[str, Any],
    *,
    field_values: dict[str, str],
    output_path: Path,
    client_name: str = "",
    project_code: str = "",
    logo_path: Path | None = None,
    on_progress: Any | None = None,
) -> tuple[Path, dict[str, str]]:
    """
    Build a deck for a package-mode proposal type.
    - Copies master/standalone template
    - Keeps only the package section slides (isolated proposal)
    - Replaces ONLY {PLACEHOLDER} tokens; raw text stays untouched
    """
    def progress(pct: int, msg: str) -> None:
        if on_progress:
            on_progress(pct, msg)

    if (pkg.get("mode") or "") != "package":
        raise ValueError(f"Tipo '{pkg.get('id')}' não é mode=package")

    progress(10, "Validando insumos e resolvendo template…")
    src = resolve_package_template(pkg)
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    progress(25, f"Copiando template {src.name}…")
    shutil.copy2(src, output_path)

    progress(40, "Identificando seção da proposta…")
    slides = resolve_package_slides(pkg, src)
    source = (pkg.get("source") or "standalone").lower()
    if source == "master" or pkg.get("section") or pkg.get("slides_1based"):
        progress(55, f"Isolando {len(slides)} slides da seção…")
        keep_only_slides(output_path, slides)

    progress(70, "Preenchendo placeholders {…}…")
    values = build_field_values(
        pkg,
        field_values=field_values,
        client_name=client_name,
        project_code=project_code,
    )
    apply_named_placeholders(output_path, values)

    # Clear leftover named tokens in the isolated deck so braces never ship
    from .engine import is_fillable_named_token, scan_named_tokens

    leftovers = {
        token: ""
        for token in scan_named_tokens(output_path)
        if token not in values and is_fillable_named_token(token)
    }
    if leftovers:
        apply_named_placeholders(output_path, leftovers)
        values.update(leftovers)

    # Optional investment table fill — total (e legado qtd/hora se ainda vierem)
    qtd = (field_values.get("quantidade") or field_values.get("pessoas") or "").strip()
    valor = (field_values.get("valor_hora") or field_values.get("valor") or "").strip()
    perfil = (field_values.get("perfil") or "").strip()
    total = (field_values.get("total") or "").strip()
    if not total and qtd and valor:
        try:
            q = float(re.sub(r"[^\d,.]", "", qtd).replace(".", "").replace(",", "."))
            v = float(re.sub(r"[^\d,.]", "", valor).replace(".", "").replace(",", "."))
            total = f"R$ {q * v:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
        except ValueError:
            total = ""

    progress(85, "Atualizando tabela de investimento…")
    if total or perfil or qtd or valor:
        fill_investment_table(
            output_path,
            perfil=perfil,
            quantidade=qtd,
            valor_hora=valor,
            total=total,
        )

    if logo_path:
        progress(92, "Aplicando logo do cliente…")
        replace_logo_cliente(output_path, Path(logo_path))

    progress(100, "Proposta pronta!")
    values["_section_slides"] = ",".join(str(s) for s in slides)
    return output_path, values
