"""PPTX template engine: parameterize, fill slots, typography, logo."""

from __future__ import annotations

import json
import re
import shutil
import sys
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt

from . import paths as P

sys.path.insert(0, str(P.APP_ROOT))
from proposal_library.importer import delete_slide  # noqa: E402
from proposal_library.placeholder_applier import (  # noqa: E402
    _iter_shapes,
    _set_shape_text,
)

# Named tokens in the master template: {TITULO_DOR}, {BREVE_DESCRICAO}, …
_NAMED_TOKEN_RE = re.compile(r"\{([A-Za-z][A-Za-z0-9_]*)\}")
# Legacy parameterized junk left in some shapes
_SLOT_JUNK_RE = re.compile(r"^S\d+_T\d+$", re.IGNORECASE)
_LOGO_TOKEN_NAMES = frozenset({"LOGO_CLIENTE", "Logo_Cliente"})
_EMPTY_TOKENS = frozenset({ "{}", "{ }" })


def shape_full_text(shape) -> str:
    return "\n".join(p.text for p in shape.text_frame.paragraphs)


def dominant_size_pt(shape) -> int | None:
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.text.strip() and run.font.size:
                return int(run.font.size.pt)
    return None


def classify_role(size: int | None, text: str, slide_no: int, slot: int) -> str:
    t = text.strip()
    if t in {"1", "2", "3", "4", "“", "R$"} or t.startswith("R$"):
        return "keep"
    if size is not None and size >= 60:
        return "display"
    if size is not None and size >= 40 and slide_no in {1, 2, 19, 20}:
        return "cover_title"
    if size is not None and size >= 28:
        return "title"
    if size == 16:
        return "subtitle"
    if size in {13, 14, 18, 21} and len(t) > 40:
        return "subtitle"
    if size is not None and size <= 14 and len(t) <= 40:
        return "label"
    return "body"


def drop_bank_slides(pptx_path: Path) -> int:
    prs = Presentation(str(pptx_path))
    removed = 0
    for i in range(len(prs.slides) - 1, -1, -1):
        if (i + 1) in P.BANK_SLIDES_1BASED:
            delete_slide(prs, i)
            removed += 1
    prs.save(str(pptx_path))
    return removed


def parameterize(pptx_path: Path) -> tuple[dict[str, str], dict[str, dict]]:
    prs = Presentation(str(pptx_path))
    originals: dict[str, str] = {}
    meta: dict[str, dict] = {}

    for si, slide in enumerate(prs.slides, start=1):
        ti = 0
        for shape in _iter_shapes(slide.shapes):
            if not hasattr(shape, "text_frame") or not shape.has_text_frame:
                continue
            text = shape_full_text(shape)
            if not text.strip() or text.strip() in P.KEEP_AS_IS:
                continue
            size = dominant_size_pt(shape)
            role = classify_role(size, text, si, ti)
            key = f"{{S{si:02d}_T{ti:02d}}}"
            originals[key] = text
            meta[key] = {"role": role, "size": size, "slide": si, "slot": ti}
            _set_shape_text(shape, key)
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = P.FONT
                    if role == "title":
                        run.font.size = Pt(P.TITLE_PT)
                    elif role == "cover_title":
                        run.font.size = Pt(P.COVER_TITLE_PT)
                    elif role == "subtitle":
                        run.font.size = Pt(P.SUBTITLE_PT)
                    elif size:
                        run.font.size = Pt(size)
            ti += 1

    prs.save(str(pptx_path))
    return originals, meta


def apply_values_and_typography(
    pptx_path: Path,
    values: dict[str, str],
    meta: dict[str, dict],
) -> int:
    ordered = sorted(values.items(), key=lambda kv: len(kv[0]), reverse=True)
    prs = Presentation(str(pptx_path))
    modified = 0

    def style_runs(shape, role: str | None, fallback_size: int | None = None):
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.font.name = P.FONT
                if role == "title":
                    run.font.size = Pt(P.TITLE_PT)
                    run.font.bold = True
                elif role == "cover_title":
                    run.font.size = Pt(P.COVER_TITLE_PT)
                    run.font.bold = True
                elif role == "subtitle":
                    run.font.size = Pt(P.SUBTITLE_PT)
                elif role == "display":
                    if run.font.size is None:
                        run.font.size = Pt(fallback_size or 64)
                elif role == "label":
                    run.font.size = Pt(min(fallback_size or 13, 14))
                elif role == "body":
                    run.font.size = Pt(fallback_size or 12)
                elif fallback_size:
                    run.font.size = Pt(fallback_size)

    for slide in prs.slides:
        for shape in _iter_shapes(slide.shapes):
            if not hasattr(shape, "text_frame"):
                continue
            full = shape_full_text(shape)
            if not full.strip():
                continue

            roles = []
            sizes = []
            for k, info in meta.items():
                if k in full:
                    roles.append(info.get("role"))
                    if info.get("size"):
                        sizes.append(info["size"])

            pending = [(k, v) for k, v in ordered if k and k in full]
            if not pending and not roles:
                style_runs(shape, None)
                continue

            remaining = []
            for old, val in pending:
                placed = False
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if old in run.text:
                            run.text = run.text.replace(old, str(val))
                            placed = True
                if not placed:
                    remaining.append((old, val))

            if remaining:
                full2 = shape_full_text(shape)
                new = full2
                for old, val in remaining:
                    new = new.replace(old, str(val))
                if new != full2:
                    _set_shape_text(shape, new)

            role = None
            for candidate in (
                "cover_title",
                "title",
                "subtitle",
                "display",
                "label",
                "body",
            ):
                if candidate in roles:
                    role = candidate
                    break
            if roles and role is None:
                role = roles[0]

            style_runs(shape, role, max(sizes) if sizes else None)
            modified += 1

    prs.save(str(pptx_path))
    return modified


def replace_logo_cliente(pptx_path: Path, logo_path: Path) -> int:
    if not logo_path or not Path(logo_path).is_file():
        return 0
    prs = Presentation(str(pptx_path))
    replaced = 0
    for slide in prs.slides:
        targets = []
        for shape in _iter_shapes(slide.shapes):
            if not hasattr(shape, "text_frame"):
                continue
            text = shape_full_text(shape)
            if "{Logo_Cliente}" in text or "{LOGO_CLIENTE}" in text:
                targets.append(shape)
        for shape in targets:
            left, top, height = shape.left, shape.top, shape.height
            el = shape._element  # noqa: SLF001
            parent = el.getparent()
            if parent is not None:
                parent.remove(el)
            slide.shapes.add_picture(str(logo_path), left, top, height=height)
            replaced += 1
    if replaced:
        prs.save(str(pptx_path))
    return replaced


def keep_only_slides(pptx_path: Path, slides_1based: list[int] | set[int]) -> int:
    """Keep only the given 1-based slide numbers; delete the rest. Returns removed count."""
    keep = {int(s) for s in slides_1based}
    prs = Presentation(str(pptx_path))
    removed = 0
    for i in range(len(prs.slides) - 1, -1, -1):
        if (i + 1) not in keep:
            delete_slide(prs, i)
            removed += 1
    prs.save(str(pptx_path))
    return removed


def _normalize_token_key(token: str) -> str:
    """Ensure key looks like {NAME}."""
    t = (token or "").strip()
    if not t:
        return ""
    if not (t.startswith("{") and t.endswith("}")):
        t = "{" + t.strip("{}") + "}"
    return t


def is_fillable_named_token(token: str) -> bool:
    """True if token should be filled by LLM / named fill (not logo, not junk, not empty)."""
    key = _normalize_token_key(token)
    if not key or key in _EMPTY_TOKENS or key in P.KEEP_AS_IS:
        return False
    inner = key[1:-1]
    if not inner or inner.isspace():
        return False
    if inner in _LOGO_TOKEN_NAMES:
        return False
    if _SLOT_JUNK_RE.match(inner):
        return False
    return bool(_NAMED_TOKEN_RE.fullmatch(key))


def scan_named_tokens(pptx_path: Path) -> dict[str, dict]:
    """
    Discover named {TOKEN} placeholders in a PPTX.
    Returns token → {slides: [1-based], count, sections: [names]}.
    Ignores empty {}, logo markers, and legacy {S##_T##} junk.
    Raw text without braces is never returned.
    """
    from .packages import read_pptx_sections

    pptx_path = Path(pptx_path)
    prs = Presentation(str(pptx_path))
    sections = read_pptx_sections(pptx_path)
    slide_to_section: dict[int, str] = {}
    for name, slides in sections.items():
        for s in slides:
            slide_to_section[s] = name

    found: dict[str, dict] = {}

    def note(token_inner: str, slide_no: int) -> None:
        key = "{" + token_inner + "}"
        if not is_fillable_named_token(key):
            return
        info = found.setdefault(
            key,
            {"slides": [], "count": 0, "sections": []},
        )
        info["count"] += 1
        if slide_no not in info["slides"]:
            info["slides"].append(slide_no)
        sec = slide_to_section.get(slide_no)
        if sec and sec not in info["sections"]:
            info["sections"].append(sec)

    def scan_text(text: str, slide_no: int) -> None:
        for m in _NAMED_TOKEN_RE.finditer(text or ""):
            note(m.group(1), slide_no)

    def walk(shapes, slide_no: int) -> None:
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                walk(shape.shapes, slide_no)
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                for row in shape.table.rows:
                    for cell in row.cells:
                        scan_text(cell.text, slide_no)
                continue
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                scan_text(shape.text_frame.text, slide_no)

    for si, slide in enumerate(prs.slides, start=1):
        walk(slide.shapes, si)

    # Stable order by first slide then name
    return dict(
        sorted(
            found.items(),
            key=lambda kv: (min(kv[1]["slides"] or [999]), kv[0]),
        )
    )


def load_named_token_catalog(pptx_path: Path | None = None) -> dict[str, dict]:
    """Catalog for the LLM: {TOKEN} → {slides, sections, count, role/writing hints}."""
    path = Path(pptx_path) if pptx_path else P.master_template_path()
    scanned = scan_named_tokens(path)
    catalog: dict[str, dict] = {}
    for key, info in scanned.items():
        sections = info.get("sections") or []
        section = sections[0] if sections else ""
        role = _guess_token_role(key, section)
        catalog[key] = {
            "role": role,
            "writing": _writing_guidance(role),
            "slide": (info.get("slides") or [None])[0],
            "slides": info.get("slides") or [],
            "sections": sections,
            "section": section,
            "count": info.get("count", 1),
            "original_example": key,
        }
    return catalog


def filter_catalog_by_section(
    catalog: dict[str, dict],
    section_name: str,
) -> dict[str, dict]:
    """Keep only tokens that appear in the given PowerPoint section."""
    wanted = (section_name or "").strip()
    if not wanted:
        return dict(catalog)
    filtered: dict[str, dict] = {}
    for key, info in catalog.items():
        sections = info.get("sections") or []
        if wanted in sections or (info.get("section") or "") == wanted:
            filtered[key] = info
    return filtered


def _guess_token_role(token: str, section: str) -> str:
    """
    Classify named tokens for writing density.

    Roles:
    - title / subtitle / cover / label / step / bullet / meta → short
    - card_desc → 1–2 sentences under a title/card
    - narrative → 2–3 short paragraphs under a section title
    """
    name = token.strip("{}").upper()
    section_u = (section or "").upper()

    if name in {
        "COD_PROJ",
        "COD_PROJETO",
        "COD_CLIENTE",
        "DATA",
        "DATA_ATUAL",
        "NOME_CLIENTE",
        "VALOR",
        "TEMPO_CONTRATO",
        "TEMPO",
        "TOTAL",
        "VALOR_ENTREGA",
        "VALOR_KICK",
        "MESES",
    } or name.startswith("VALOR_"):
        return "meta"
    if name.startswith("STEP_") or name == "SEMANAS":
        return "step"
    if name.startswith("BULLET") or name.endswith("_BULLETS") or name.startswith("ITENS_"):
        return "bullet"
    if name.startswith("SUB_") or name.startswith("SUBTITULO") or name.startswith("BREVE"):
        return "subtitle"
    if name == "DESC_TITULO":
        return "subtitle"
    if (
        name.startswith("TITULO")
        or name.endswith("_TITULO")
        or name in {"OBJETIVO_TITULO", "RESULT_TITULO", "ENTREGA_TITULO"}
    ):
        return "title"
    if name.startswith("CARD_") and name.endswith("_TITULO"):
        return "label"
    if name.startswith("EXEC_") and name.endswith("_TITULO"):
        return "label"
    if name.startswith("CARD_") and not name.endswith("_DESC"):
        return "label"

    # Supporting text under a card / item / pillar (1–2 sentences)
    if (
        name.startswith("CARD_")
        or re.fullmatch(r"DESC_[1-7]_CARD", name)
        or re.fullmatch(r"DESC_[1-7]_OPORT", name)
        or re.fullmatch(r"DESC_[1-7]_ENTREGAVEL", name)
        or re.fullmatch(r"DESC_SUB_DOR_[1-7]", name)
        or name.endswith("_DESAFIO_DESC")
        or (name.startswith("EXEC_") and name.endswith("_DESC"))
        or name.startswith("IMPACTO_")
        or name.startswith("PREM_")
        or name.startswith("RESTR_")
        or name.startswith("REST_")
        or name.startswith("PREMISSA_")
        or name.startswith("RESTRICAO_")
        or (name.startswith("ITEM_") and name.endswith("_DESC"))
        or (name.startswith("ITEM_") and name.endswith("_OPORTUNIDADE"))
    ):
        return "card_desc"

    if name.endswith("_CARD"):
        return "label"
    if name.endswith("_ITEM"):
        return "bullet"

    if name.startswith("ITEM_"):
        return "label"

    # Main narrative body under the slide title (2–3 paragraphs)
    if (
        name.endswith("_DESC")
        or name.startswith("DESC_")
        or "DESCRICAO" in name
        or "DECRICAO" in name
        or name.startswith("RESULT_")
        or name.startswith("ENTREGA_")
        or name.startswith("DESC_SUB_")
    ):
        return "narrative"

    if section_u.startswith("CAPA"):
        return "cover"
    return "narrative"


def _writing_guidance(role: str) -> str:
    return {
        "title": "4–10 palavras; impacto; sem parágrafo",
        "subtitle": "1–2 linhas; complementar o título; sem parágrafo longo",
        "cover": "texto curto de capa; sem feature/arquitetura",
        "label": "rótulo curto (2–5 palavras)",
        "step": "etapa/indicador curto",
        "bullet": "lista objetiva; sem parágrafo",
        "meta": "metadado curto (código/data/cliente)",
        "card_desc": "1–2 frases (10–22 palavras) sob o título do card",
        "narrative": "2–3 parágrafos curtos separados por \\n\\n (45–90 palavras)",
    }.get(role, "texto executivo objetivo")


# Sections isolated into package modes — excluded from Livre decks
_LIVRE_EXCLUDED_SECTIONS = (
    "Professional Service",
    "SUPORTE",
    "CONTROLE DE ACESSO (PASSLOG)",
    "DISCOVERY",
    "CLARION",
    "Escopo Fechado (DP World)",
)


def livre_slide_indices(pptx_path: Path | None = None) -> list[int]:
    """
    1-based slides for Livre mode: all sections except package-only ones
    (Professional Service, SUPORTE, PassLog, Discovery, Clarion,
    Escopo Fechado). If sections missing, keep the whole deck.
    """
    from .packages import read_pptx_sections

    path = Path(pptx_path) if pptx_path else P.master_template_path()
    sections = read_pptx_sections(path)
    prs = Presentation(str(path))
    all_slides = list(range(1, len(prs.slides) + 1))
    if not sections:
        return all_slides
    excluded: set[int] = set()
    for name in _LIVRE_EXCLUDED_SECTIONS:
        excluded.update(sections.get(name) or [])
    keep = [s for s in all_slides if s not in excluded]
    return keep or all_slides


def build_livre_deck(
    values: dict[str, str],
    *,
    output_path: Path,
    logo_path: Path | None = None,
    client_name: str = "",
    project_code: str = "",
    exclude_professional_service: bool = True,
) -> Path:
    """
    Copy slide-mestre, optionally drop package-only sections
    (Professional Service, SUPORTE, PassLog, Discovery, Clarion,
    Escopo Fechado), replace ONLY named {TOKENS}, apply logo.
    Raw text stays untouched.
    """
    src = P.master_template_path()
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(src, output_path)

    if exclude_professional_service:
        keep = livre_slide_indices(src)
        keep_only_slides(output_path, keep)

    today = date.today().strftime("%d/%m/%Y")
    merged: dict[str, str] = {
        "{COD_PROJ}": (project_code or "").strip() or "A definir",
        "{COD_CLIENTE}": (project_code or "").strip() or "A definir",
        "{DATA}": today,
        "{DATA_ATUAL}": today,
        "{NOME_CLIENTE}": (client_name or "").strip() or "Cliente",
        "{Nome_cliente}": (client_name or "").strip() or "Cliente",
    }
    for k, v in values.items():
        if k.startswith("_"):
            continue
        key = _normalize_token_key(k)
        if not key or not is_fillable_named_token(key):
            continue
        # Empty string is intentional (cronograma/premissas sem dado)
        merged[key] = "" if v is None else str(v)

    apply_named_placeholders(output_path, merged)
    clear_legacy_slot_tokens(output_path)
    if logo_path:
        replace_logo_cliente(output_path, Path(logo_path))
    return output_path


def apply_named_placeholders(
    pptx_path: Path,
    values: dict[str, str],
    *,
    force: bool = False,
) -> int:
    """
    Replace named placeholders like {COD_CLIENTE} in text frames and tables.
    Skips empty {} and never invents content for non-placeholder text.
    force=True also replaces logo markers (used to clear them when unused).
    """
    if not values:
        return 0
    # Ignore empty / logo markers (logo handled by replace_logo_cliente)
    clean: dict[str, str] = {}
    for k, v in values.items():
        key = _normalize_token_key(k)
        if not key or key in _EMPTY_TOKENS or key.startswith("_"):
            continue
        inner = key[1:-1]
        if not force and (
            key in P.KEEP_AS_IS or inner in _LOGO_TOKEN_NAMES
        ):
            continue
        if _SLOT_JUNK_RE.match(inner):
            continue
        clean[key] = "" if v is None else str(v)

    # Alias Logo_Cliente ↔ LOGO_CLIENTE for any residual text fill
    if "{LOGO_CLIENTE}" in clean and "{Logo_Cliente}" not in clean:
        clean["{Logo_Cliente}"] = clean["{LOGO_CLIENTE}"]
    if "{Logo_Cliente}" in clean and "{LOGO_CLIENTE}" not in clean:
        clean["{LOGO_CLIENTE}"] = clean["{Logo_Cliente}"]

    if not clean:
        return 0
    ordered = sorted(clean.items(), key=lambda kv: len(kv[0]), reverse=True)
    prs = Presentation(str(pptx_path))
    modified = 0

    def replace_in_text_frame(tf) -> bool:
        full = "\n".join(p.text for p in tf.paragraphs)
        if not any(k in full for k, _ in ordered):
            return False
        remaining = []
        for old, val in ordered:
            placed = False
            for para in tf.paragraphs:
                for run in para.runs:
                    if old in run.text:
                        run.text = run.text.replace(old, val)
                        placed = True
            if not placed and old in "\n".join(p.text for p in tf.paragraphs):
                remaining.append((old, val))
        if remaining:
            full2 = "\n".join(p.text for p in tf.paragraphs)
            new = full2
            for old, val in remaining:
                new = new.replace(old, val)
            if new != full2:
                _set_shape_text_frame(tf, new)
        final_text = "\n".join(p.text for p in tf.paragraphs)
        if "\n\n" in final_text or len(final_text) > 120:
            # Narrative body fields may contain 2–3 short paragraphs. Let
            # PowerPoint reduce the font only when needed instead of clipping.
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        return True

    def walk(shapes):
        nonlocal modified
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                walk(shape.shapes)
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if replace_in_text_frame(cell.text_frame):
                            modified += 1
                continue
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                if replace_in_text_frame(shape.text_frame):
                    modified += 1

    for slide in prs.slides:
        walk(slide.shapes)

    if modified:
        prs.save(str(pptx_path))
    return modified


def clear_legacy_slot_tokens(pptx_path: Path) -> int:
    """Remove residual {S##_T##} placeholders without touching other raw text."""
    prs = Presentation(str(pptx_path))
    modified = 0

    def clear_text_frame(tf) -> bool:
        changed = False
        for para in tf.paragraphs:
            for run in para.runs:
                new = re.sub(r"\{S\d+_T\d+\}", "", run.text, flags=re.IGNORECASE)
                if new != run.text:
                    run.text = new
                    changed = True
        if changed:
            return True

        full = "\n".join(p.text for p in tf.paragraphs)
        new = re.sub(r"\{S\d+_T\d+\}", "", full, flags=re.IGNORECASE)
        if new != full:
            _set_shape_text_frame(tf, new)
            return True
        return False

    def walk(shapes) -> None:
        nonlocal modified
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                walk(shape.shapes)
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if clear_text_frame(cell.text_frame):
                            modified += 1
                continue
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                if clear_text_frame(shape.text_frame):
                    modified += 1

    for slide in prs.slides:
        walk(slide.shapes)

    if modified:
        prs.save(str(pptx_path))
    return modified


def _set_shape_text_frame(tf, text: str) -> None:
    """Write plain text into a text frame preserving first paragraph style lightly."""
    lines = text.split("\n")
    for i, para in enumerate(tf.paragraphs):
        if i < len(lines):
            if para.runs:
                para.runs[0].text = lines[i]
                for run in para.runs[1:]:
                    run.text = ""
            else:
                para.text = lines[i]
        else:
            if para.runs:
                for run in para.runs:
                    run.text = ""
            else:
                para.text = ""
    # extra lines
    first = tf.paragraphs[0] if tf.paragraphs else None
    for line in lines[len(tf.paragraphs) :]:
        p = tf.add_paragraph()
        p.text = line
        if first is not None and first.runs:
            # style copy is best-effort
            pass


def ensure_parameterized_template(force: bool = False) -> tuple[Path, dict, dict]:
    """
    Ensure variaveis PPTX + slots JSON exist.
    Returns (template_path, originals, meta).
    """
    tpl = P.template_vars_path()
    slots_file = P.slots_path()

    if (
        not force
        and tpl.is_file()
        and slots_file.is_file()
    ):
        data = json.loads(slots_file.read_text(encoding="utf-8"))
        return tpl, data["originals"], data["meta"]

    model = P.resolve_model()
    P.UNISANTA.mkdir(parents=True, exist_ok=True)
    shutil.copy2(model, tpl)
    drop_bank_slides(tpl)
    originals, meta = parameterize(tpl)
    slots_file.write_text(
        json.dumps({"meta": meta, "originals": originals}, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    return tpl, originals, meta


def build_deck(
    values: dict[str, str],
    *,
    output_path: Path,
    logo_path: Path | None = None,
    force_reparam: bool = False,
) -> Path:
    """Copy parameterized template, fill values, apply logo, save."""
    tpl, originals, meta = ensure_parameterized_template(force=force_reparam)

    # Fill missing keys with originals (keeps layout text if LLM skipped a slot)
    merged = dict(originals)
    for k, v in values.items():
        if v is not None and str(v).strip():
            merged[k] = str(v)

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(tpl, output_path)
    apply_values_and_typography(output_path, merged, meta)
    if logo_path:
        replace_logo_cliente(output_path, Path(logo_path))
    return output_path


def load_slot_catalog() -> dict:
    """Catalog for the LLM: key -> {role, original preview}."""
    _, originals, meta = ensure_parameterized_template()
    catalog = {}
    for key, original in originals.items():
        info = meta.get(key, {})
        catalog[key] = {
            "role": info.get("role", "body"),
            "slide": info.get("slide"),
            "original_example": original[:280],
        }
    return catalog
